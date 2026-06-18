from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# -------------------- Цветовая палитра --------------------
PRIMARY_BLUE = RGBColor(26, 60, 110)
LIGHT_BLUE = RGBColor(230, 240, 250)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(200, 200, 200)
DARK_GRAY = RGBColor(80, 80, 80)
RED = RGBColor(200, 50, 50)
GREEN = RGBColor(50, 180, 50)

# -------------------- Вспомогательные функции --------------------
def set_text_frame_props(tf, font_name='Roboto', size=16, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tf.font.name = font_name
    tf.font.size = Pt(size)
    tf.font.bold = bold
    tf.font.color.rgb = color
    tf.paragraphs[0].alignment = align
    # для списков можно добавить отступы, но это упростим

def add_textbox(slide, left, top, width, height, text, **kwargs):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.text = text
    set_text_frame_props(tf, **kwargs)
    return shape

def add_shape_with_text(slide, left, top, width, height, text, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, fill_color=LIGHT_BLUE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=14, bold=False):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Roboto'
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = font_color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_progress_bar(slide, left, top, width, height, percentage, fill_color=PRIMARY_BLUE, bg_color=GRAY):
    # фоновая полоса
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    # залитая часть
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width * percentage / 100), Inches(height))
    fill.fill.solid()
    fill.fill.fore_color.rgb = fill_color
    fill.line.fill.background()
    # текст процента
    add_textbox(slide, left + width * 0.4, top - 0.1, 1.5, 0.5, f"{percentage}%", size=18, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)

# -------------------- Создание презентации --------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Макет для заголовка и контента (используем встроенный)
title_slide_layout = prs.slide_layouts[0]   # Заголовок + подзаголовок
blank_slide_layout = prs.slide_layouts[6]   # Пустой

# 1. Титульный слайд
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "ГосБанк"
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
subtitle = slide.placeholders[1]
subtitle.text = "Финансовая платформа\nРазработчик: Ярцев В.А | Группа: 24290907/2093\nGitHub | Сайт"

# 2. Назначение, предметная область, защита
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Назначение и предметная область", size=28, bold=True, color=PRIMARY_BLUE)
# три колонки
lefts = [1, 5, 9]
for i, (title_text, desc) in enumerate([
    ("Одобрение карт, кредиты, овердрафт", "Оформление карт, выдача кредитов и овердрафта"),
    ("Вход и верификация через Госуслуги", "Безопасная аутентификация через Госуслуги"),
    ("Администрирование", "Одобрение новых администраторов")
]):
    add_shape_with_text(slide, lefts[i], 1.5, 2.5, 1.0, title_text, font_size=16, bold=True, fill_color=LIGHT_BLUE, line_color=PRIMARY_BLUE, font_color=PRIMARY_BLUE)
    add_textbox(slide, lefts[i], 2.7, 2.5, 0.8, desc, size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
# иконки (символы) можно добавить над карточками, но пропустим для простоты

# 3. Срок, удалённость, верификация
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Срок, удалённость и верификация", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 3, 1.2, "Срок разработки: 5 месяцев", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK)
add_progress_bar(slide, 1, 2.8, 3, 0.3, 100, fill_color=ORANGE) # прогресс не имеет смысла, но покажем
add_shape_with_text(slide, 5, 1.5, 3, 1.2, "☁️ Удалённость серверов", fill_color=WHITE, line_color=PRIMARY_BLUE)
add_shape_with_text(slide, 9, 1.5, 3, 1.2, "✅ Верификация только через Госуслуги", fill_color=WHITE, line_color=PRIMARY_BLUE)

# 4. Три столпа (Доступность, Конфиденциальность, Целостность)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Ключевые принципы безопасности", size=28, bold=True, color=PRIMARY_BLUE)
for i, word in enumerate(["Доступность", "Конфиденциальность", "Целостность"]):
    add_shape_with_text(slide, 1 + i*4, 1.5, 3, 2.0, word, fill_color=LIGHT_BLUE, line_color=PRIMARY_BLUE, font_color=PRIMARY_BLUE, font_size=24, bold=True)

# 5. Нагрузка и аудитория
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Нагрузка и аудитория", size=28, bold=True, color=PRIMARY_BLUE)
add_textbox(slide, 2, 2, 4, 1.5, "До 1 млн пользователей", size=48, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, 2, 3.5, 4, 0.8, "только граждане РФ", size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 6. Модель ЖЦ (спиральная) + причины
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Модель жизненного цикла: Спиральная", size=28, bold=True, color=PRIMARY_BLUE)
# рисуем схему спирали (упрощённо – четыре круга)
for i in range(4):
    left = 1 + i * 2.5
    top = 2.5 - i * 0.3
    add_shape_with_text(slide, left, top, 2, 2, f"Итерация {i+1}", shape_type=MSO_SHAPE.OVAL, fill_color=LIGHT_BLUE, line_color=PRIMARY_BLUE, font_color=PRIMARY_BLUE, font_size=14)
# причины в столбик
reasons = ["Позволяет управлять рисками", "Обеспечивает гибкость", "Раннее представление продукта", "Снижает риск провала"]
for idx, reason in enumerate(reasons):
    add_textbox(slide, 8, 1.5 + idx * 0.6, 4, 0.5, f"✓ {reason}", size=16, color=BLACK)

# 7. Итерации (1 и 2)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Итерации разработки", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 4.5, 2, "Итерация 1: Базовый каркас\nАрхитектура, БД, API", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK)
add_shape_with_text(slide, 7, 1.5, 4.5, 2, "Итерация 2: Аутентификация и безопасность\nГосуслуги, роли, шифрование", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK)
# стрелка между
add_textbox(slide, 5.8, 2, 0.8, 0.8, "→", size=36, color=ORANGE, align=PP_ALIGN.CENTER)

# 8. Уровни доступа (пирамида)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Уровни доступа", size=28, bold=True, color=PRIMARY_BLUE)
levels = ["Суперадмин", "Админ", "Зарегистрированный пользователь", "Гость"]
heights = [0.8, 0.8, 0.8, 0.8]
top_start = 1.8
for i, level in enumerate(levels):
    y = top_start + i * 1.0
    width = 8 - i * 1.2
    left = (12 - width) / 2
    add_shape_with_text(slide, left, y, width, 0.7, level, fill_color=LIGHT_BLUE if i<2 else WHITE, line_color=PRIMARY_BLUE, font_color=PRIMARY_BLUE, font_size=18, bold=True)

# 9. Путь пользователя (блок-схема)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Путь пользователя", size=28, bold=True, color=PRIMARY_BLUE)
steps = ["Главная", "Калькулятор овердрафта", "Подать заявку", "Регистрация", "Вход через Госуслуги", "Ожидание подтверждения"]
for i, step in enumerate(steps):
    left = 0.5 + i * 1.9
    add_shape_with_text(slide, left, 2.5, 1.7, 0.8, step, fill_color=LIGHT_BLUE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=12)
    if i < len(steps)-1:
        add_textbox(slide, left + 1.7, 2.4, 0.3, 0.8, "→", size=24, color=ORANGE, align=PP_ALIGN.CENTER)

# 10. Интерфейс (макет, видео, гифка) – просто заглушки
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Интерфейс и медиа", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 4, 2.5, "Макет\n(скриншот)", fill_color=GRAY, line_color=DARK_GRAY, font_color=WHITE, font_size=20)
add_shape_with_text(slide, 6, 1.5, 3, 2.5, "Видео\n▶️", fill_color=GRAY, line_color=DARK_GRAY, font_color=WHITE, font_size=20)
add_shape_with_text(slide, 9.5, 1.5, 2.5, 2.5, "GIF\n🔄", fill_color=GRAY, line_color=DARK_GRAY, font_color=WHITE, font_size=20)

# 11. Принципы удобного GUI
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Принципы удобного GUI", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 4.5, 1.5, "Принцип простоты\nРегистрация – в минимум кликов", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=16)
add_shape_with_text(slide, 7, 1.5, 4.5, 1.5, "Принцип видимости\nВсе нужные данные видны при решении задачи", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=16)

# 12. ГОСТы и руководство оператора
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "ГОСТы и руководство оператора", size=28, bold=True, color=PRIMARY_BLUE)
add_textbox(slide, 1, 1.5, 4, 2, "ГОСТы:\n19.104-78\n19.106-78\n19.505-79\nISO 9126", size=16, color=BLACK)
add_textbox(slide, 6, 1.5, 5, 2, "Руководство оператора (пример ошибки)", size=16, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 6, 2.5, 5, 1.2, "Ошибка: Неверный логин или пароль", fill_color=RED, line_color=RED, font_color=WHITE, font_size=18)
add_textbox(slide, 6, 3.8, 5, 1, "→ Проверить данные\n→ При повторении обратиться к администратору", size=14, color=DARK_GRAY)

# 13. Измеримые критерии качества
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Критерии качества", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 5, 2.5, "Надёжность\n✅ Сохранение формы при обрыве\n✅ Wipe and Rebuild", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=16)
add_shape_with_text(slide, 7, 1.5, 5, 2.5, "Удобство использования\n✅ Автофокус\n✅ Подсказки в полях\n✅ Сообщения об ошибках", fill_color=WHITE, line_color=PRIMARY_BLUE, font_color=BLACK, font_size=16)

# 14. Методы и техники тестирования
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Методы и техники тестирования", size=28, bold=True, color=PRIMARY_BLUE)
add_textbox(slide, 1, 1.5, 5, 0.8, "Техники:", size=20, bold=True, color=PRIMARY_BLUE)
techniques = ["EP (эквивалентное разделение)", "BVA (анализ граничных значений)", "EG (предугадывание ошибки)"]
for i, t in enumerate(techniques):
    add_textbox(slide, 1.5, 2.3 + i*0.5, 4, 0.4, f"• {t}", size=16, color=BLACK)
add_textbox(slide, 7, 1.5, 5, 0.8, "Методы:", size=20, bold=True, color=PRIMARY_BLUE)
methods = ["Black Box", "Manual Testing", "Динамическое тестирование"]
for i, m in enumerate(methods):
    add_textbox(slide, 7.5, 2.3 + i*0.5, 4, 0.4, f"• {m}", size=16, color=BLACK)

# 15. Виды и типы тестирования
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Виды и типы тестирования", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.8, 4, 1.5, "Позитивные тест-кейсы", fill_color=GREEN, line_color=GREEN, font_color=WHITE, font_size=20)
add_shape_with_text(slide, 7, 1.8, 4, 1.5, "Негативные тест-кейсы", fill_color=RED, line_color=RED, font_color=WHITE, font_size=20)
add_textbox(slide, 1, 3.8, 10, 0.5, "Тип: Функциональное тестирование пользовательского GUI", size=18, color=BLACK, align=PP_ALIGN.CENTER)

# 16. Конкретные тест-кейсы (№1 и №4)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Примеры тест-кейсов", size=28, bold=True, color=PRIMARY_BLUE)
add_shape_with_text(slide, 1, 1.5, 5, 2, "Позитивный тест-кейс №1\n(успешная регистрация)", fill_color=WHITE, line_color=GREEN, font_color=BLACK, font_size=18)
add_shape_with_text(slide, 7, 1.5, 5, 2, "Негативный тест-кейс №4\n(ошибка валидации)", fill_color=WHITE, line_color=RED, font_color=BLACK, font_size=18)

# 17. Тестовое покрытие (51,6%)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Тестовое покрытие", size=28, bold=True, color=PRIMARY_BLUE)
add_textbox(slide, 2, 2, 8, 1.5, "51.6%", size=72, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_progress_bar(slide, 2, 3.5, 8, 0.6, 51.6, fill_color=ORANGE, bg_color=GRAY)
add_textbox(slide, 2, 4.3, 8, 0.5, "Покрыто 16 из 31 проверок", size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 18. Вывод и перспективы (выполнено 45%)
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 0.3, 12, 0.8, "Вывод и перспективы", size=28, bold=True, color=PRIMARY_BLUE)
add_progress_bar(slide, 2, 1.5, 8, 0.6, 45, fill_color=PRIMARY_BLUE, bg_color=GRAY)
add_textbox(slide, 2, 2.5, 8, 0.4, "Проект реализован на 45%", size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# три колонки
columns = [
    ("Что удалось", ["Концептуальная часть", "Интерфейс", "Тестирование", "Документирование"]),
    ("Что хотелось бы изменить", ["Организационные аспекты", "Технические решения"]),
    ("Перспективы", ["Государственное финансирование", "Масштабируемость"])
]
for idx, (title, items) in enumerate(columns):
    left = 1 + idx * 3.8
    add_textbox(slide, left, 3.5, 3.5, 0.5, title, size=20, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)
    y = 4.0
    for item in items:
        add_textbox(slide, left + 0.2, y, 3, 0.4, f"• {item}", size=14, color=BLACK)
        y += 0.4

# 19. Заключительный слайд
slide = prs.slides.add_slide(blank_slide_layout)
add_textbox(slide, 0.5, 2, 12, 2, "Спасибо за внимание!", size=48, bold=True, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 4.5, 12, 0.8, "Готов ответить на ваши вопросы", size=24, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 6, 12, 0.6, "GitHub | Сайт", size=20, color=PRIMARY_BLUE, align=PP_ALIGN.CENTER)

# Сохраняем
prs.save("GosBank_improved.pptx")
print("Презентация успешно создана: GosBank_improved.pptx")